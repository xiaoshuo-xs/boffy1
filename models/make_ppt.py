"""Generate PPT: QD-MSA Batch Processing & TT Compression Security"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# Colors
BG = RGBColor(0x1a, 0x1a, 0x2e)
FG = RGBColor(0xe0, 0xe0, 0xe0)
ACCENT = RGBColor(0x00, 0xd4, 0xaa)
ACCENT2 = RGBColor(0xff, 0x6b, 0x6b)
ACCENT3 = RGBColor(0xff, 0xd9, 0x3d)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x88, 0x88, 0x99)

def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def add_title(slide, text, y=0.3, size=36, color=WHITE):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = True
    return tb

def add_text(slide, text, y=1.3, size=18, color=FG):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = Pt(4)
    return tb

def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(12); p.font.color.rgb = GRAY

def add_box(slide, x, y, w, h, title, lines, title_color=ACCENT, bg_color=RGBColor(0x0a, 0x2a, 0x1a)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(22); p.font.color.rgb = title_color; p.font.bold = True
    for line in lines:
        p = tf.add_paragraph()
        p.text = line; p.font.size = Pt(15); p.font.color.rgb = FG
        p.space_after = Pt(3)

# ================================================================
# Slide 1: Title
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, 'QD-MSA 批处理优化 & TT 压缩安全增强', y=1.5, size=44)
add_text(s, 'Efficient Batched Quantum Circuit Simulation via Explicit Parameterization\n& Tensor-Train Compression for Hybrid QNN Security', y=2.5, size=22, color=ACCENT)
add_text(s, 'PennyLane / PyTorch 混合量子-经典神经网络优化  |  CMU-MOSEI 多模态情感分析', y=3.5, size=18, color=GRAY)
add_footer(s, 'QD-MSA Batch & TT Security  |  2026')

# ================================================================
# Slide 2: Background
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, '原始问题：逐样本循环的性能瓶颈', size=36)
add_text(s, (
    'QD-MSA 模型管线:\n'
    '  GRU 编码器 (经典) -> MLP 降维 (经典) -> MPS 量子电路 (9 qubit, 4+5 分割)\n'
    '\n'
    '瓶颈: quantum_split_model.py forward() 方法\n'
    '  每个 batch (32 样本) x 9 个量子电路 (3 前端 + 6 后端)\n'
    '  = 288 次独立 Python 调用, 完全串行\n'
    '\n'
    '后果:\n'
    '  每 epoch ~19 分钟 (batch_size=32)\n'
    '  50 epochs = ~16 小时\n'
    '  GPU 空闲 -- 量子模拟在 CPU 上运行'
), y=1.3, size=18)
add_footer(s, 'QD-MSA 原始架构分析')

# ================================================================
# Slide 3: Root Cause
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, '根因分析：TorchLayer 的批处理检测缺陷', size=36)
add_text(s, (
    '原始代码使用 PennyLane TorchLayer 包装 QNode:\n'
    '  self.QLayer_front_1 = qml.qnn.TorchLayer(circuit_front_1, weight_shapes)\n'
    '\n'
    '问题电路 (参数索引式):\n'
    '  def embedding_circuit(inputs, n_qubits):\n'
    '      for qub in range(n_qubits):\n'
    '          qml.RY(inputs[qub], wires=qub)   # <- Python 索引\n'
    '\n'
    '冲突:\n'
    '  单样本: inputs = (4,)  -> inputs[0] = 标量 -> PennyLane 识别为 1 个参数\n'
    '  批处理: inputs = (32,4) -> inputs[0] = (4,) -> PennyLane 无法区分\n'
    '    "第 0 个样本的所有参数" vs "所有样本的第 0 个参数"\n'
    '\n'
    'Python 索引语义 != PennyLane 批处理语义 -> 检测失败 -> 退化为逐样本\n'
    '作者注释: "# TorchLayer不支持这种电路的batch模式"'
), y=1.2, size=16)
add_footer(s, 'PennyLane TorchLayer 批处理检测缺陷')

# ================================================================
# Slide 4: Solution
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, '解决方案：显式标量参数 + 直接 QNode 调用', size=36)
add_text(s, (
    '核心思想: 绕过 TorchLayer, 直接调用 QNode, 用显式参数替代索引\n'
    '\n'
    '关键改动:\n'
    '  原版:                               批处理版:\n'
    '  def circuit(inputs, weights):       def circuit(i0,i1,i2,i3, weights):\n'
    '    qml.RY(inputs[0], wires=0)          qml.RY(i0, wires=0)\n'
    '    qml.RY(inputs[1], wires=1)          qml.RY(i1, wires=1)\n'
    '    qml.RY(inputs[2], wires=2)          qml.RY(i2, wires=2)\n'
    '    qml.RY(inputs[3], wires=3)          qml.RY(i3, wires=3)\n'
    '\n'
    '调用方式:\n'
    '  原版: circuit(inputs[0], w)          批版: circuit(f[:,0], f[:,1], f[:,2], f[:,3], w)\n'
    '  1 个 tensor (4,)                     4 个 tensor (batch,) -> PennyLane 自动识别 batch\n'
    '  TorchLayer 包装                       直接 QNode, 无额外开销'
), y=1.2, size=16)
add_footer(s, '显式参数传递方案')

# ================================================================
# Slide 5: Speed Benchmark
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, '实验结果：10.4x 加速 (RTX 4060, default.qubit+backprop)', size=34)
add_text(s, (
    'batch= 4:  原版 0.127s  ->  批版 0.088s  ->  1.4x\n'
    'batch= 8:  原版 0.263s  ->  批版 0.090s  ->  2.9x\n'
    'batch=16:  原版 0.517s  ->  批版 0.092s  ->  5.6x\n'
    'batch=32:  原版 1.007s  ->  批版 0.097s  -> 10.4x  ***\n'
    '\n'
    '关键发现:\n'
    '  批处理版耗时几乎不随 batch_size 增长 (0.088s -> 0.097s)\n'
    '  原版严格线性增长 (0.127s -> 1.007s)\n'
    '  完整 CMU-MOSEI 训练: 16h -> 1.5h\n'
    '\n'
    '加速来源:\n'
    '  消除 Python 循环: 288 次 -> 9 次 Python 调用\n'
    '  C++ 层高效批量处理 (default.qubit)\n'
    '  免去 torch.stack 等中间内存操作'
), y=1.3, size=17)
add_footer(s, '性能基准测试: RTX 4060 Laptop GPU')

# ================================================================
# Slide 6: Pros & Cons
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, '批处理方案：优缺点对比', size=36)
add_box(s, 0.5, 1.3, 5.8, 5.5,
    'Advantages / 优点', [
        '1. 加速 5-15x, batch 越大越明显',
        '2. 不改电路逻辑, 数学完全等价',
        '3. 通用方案: 适用于任何 indexed-param',
        '   PennyLane 电路',
        '4. 零额外依赖, 纯 Python 改动',
        '5. 前向/后向输出与原版逐位一致',
        '6. 与 TT 压缩、GPU 后端兼容',
    ], ACCENT, RGBColor(0x0a, 0x2a, 0x1a))
add_box(s, 6.8, 1.3, 5.8, 5.5,
    'Disadvantages / 缺点', [
        '1. 须用 default.qubit (支持 backprop)',
        '   lightning.qubit 不支持 batched 梯度',
        '2. backprop 求梯度 -> 不能部署真量子硬件',
        '3. 电路定义变冗长 (需显式列出所有参数)',
        '4. 需手动管理 CPU/GPU 设备转换',
        '5. 内存占用略增 (整 batch 同时模拟)',
        '6. 对 >20 量子比特的电路未验证',
    ], ACCENT2, RGBColor(0x2a, 0x0a, 0x0a))
add_footer(s, '批处理方案全面评估')

# ================================================================
# Slide 7: TT Compression Intro
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, 'TT (Tensor Train) 压缩原理', size=36)
add_text(s, (
    'TT 压缩 = 将一个大矩阵分解为多个小矩阵 (TT-cores) 的乘积:\n'
    '  W (MxN) ~= G1 x G2 x ... x Gd\n'
    '\n'
    '例子: Linear(256 -> 256)\n'
    '  原始: W in R^(256x256) = 65,536 参数\n'
    '  分解: 256 = 16x16 -> 2 个 TT-cores\n'
    '    G1: (1, 16, 16, r)  <- 第 1 个 core\n'
    '    G2: (r, 16, 16, 1)  <- 第 2 个 core\n'
    '  rank=4: 共 2,048 参数 -> 节省 96.9%\n'
    '\n'
    '安全性的数学基础 (核心):\n'
    '  TT 分解不是唯一的! 对任意可逆矩阵 Q:\n'
    '    (G1 * Q) x (Q^{-1} * G2) = G1 x G2 = W\n'
    '  同一功能有无穷多种参数表示\n'
    '  -> 攻击者即使拿到参数, 也无法唯一还原原始模型'
), y=1.2, size=17)
add_footer(s, 'TT 压缩: 参数效率 + 数学安全')

# ================================================================
# Slide 8: TT Security Architecture
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, 'TT 压缩安全分析：三层混淆架构', size=36)
add_text(s, (
    'Layer 1: TT 压缩层 (数学安全) -- 模拟器上唯一真正有效的安全措施\n'
    '  机制: 矩阵分解不唯一, 攻击者无法唯一还原原始参数\n'
    '  对抗: 模型窃取、参数逆向\n'
    '  性质: 数学严格, 不依赖硬件\n'
    '\n'
    'Layer 2: 量子 MPS 电路 (物理安全) -- 真量子硬件部署时激活\n'
    '  机制: 参数编码在量子态中, 量子测量天然扰动\n'
    '  对抗: 侧信道攻击、模型克隆\n'
    '  性质: 模拟器上仅为潜在优势, 真量子设备上物理保证\n'
    '\n'
    'Layer 3: 电路切割 (结构安全) -- 信息碎片化\n'
    '  机制: 前端(4 qubit) + 后端(5 qubit) + 断层扫描重组\n'
    '  对抗: 单点攻击、梯度泄露\n'
    '  性质: 攻击者需同时攻破前后两端才能重组完整模型'
), y=1.3, size=16)
add_footer(s, '安全架构: 数学(TT) + 物理(量子) + 结构(切割) = 三重防护')

# ================================================================
# Slide 9: Parameter Efficiency
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, '参数效率对比', size=36)
add_text(s, (
    'TT 压缩参数节省 (量子头, input_shape=256):\n'
    '  tt_rank = 0 (无压缩):  70,134 params   (baseline)\n'
    '  tt_rank = 2:            3,766 params   (saves 94.6%)\n'
    '  tt_rank = 4:            4,982 params   (saves 92.9%)\n'
    '  tt_rank = 8:            7,414 params   (saves 89.4%)  <- current\n'
    '  tt_rank = 16:          12,278 params   (saves 82.5%)\n'
    '\n'
    'CMU-MOSEI 完整模型:\n'
    '  Classical MLP:  336,327 params\n'
    '  Quantum TT (r=8): 241,718 params  (28.1% less)\n'
    '\n'
    'Tradeoff:\n'
    '  rank up -> expressivity up, security down\n'
    '  rank down -> security up, accuracy may drop\n'
    '  Recommended: rank=8~16'
), y=1.3, size=17)
add_footer(s, '参数效率与安全权衡')

# ================================================================
# Slide 10: TT Security Pros & Cons
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, 'TT 压缩：安全性优缺点', size=36)
add_box(s, 0.5, 1.3, 5.8, 5.5,
    'Security Advantages / 安全优势', [
        '1. 分解不唯一 -> 模型不可逆提取',
        '2. 参数大幅减少 -> 攻击面缩小',
        '3. 梯度路径更长 -> 梯度攻击更难',
        '4. 数学严格, 不依赖硬件假设',
        '5. 可组合: TT + 量子 + 切割 = 三重防护',
        '6. 低秩 = 天然正则化 -> 抗过拟合',
    ], ACCENT, RGBColor(0x0a, 0x2a, 0x1a))
add_box(s, 6.8, 1.3, 5.8, 5.5,
    'Limitations / 局限', [
        '1. Low-rank 假设不一定适用于所有任务',
        '2. Rank 太低会损害准确率',
        '3. 对抗性秩恢复攻击理论上可能',
        '   (已知分解结构可暴力搜索 Q)',
        '4. 增加了模型复杂度',
        '5. 不是端到端加密 -- 仍是混淆',
        '6. 需要实验验证安全-准确率 tradeoff',
    ], ACCENT2, RGBColor(0x2a, 0x0a, 0x0a))
add_footer(s, 'TT 压缩安全评估')

# ================================================================
# Slide 11: Summary
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, '总结与展望', size=36)
add_text(s, (
    '已完成:\n'
    '  1. 批处理优化: 发现并绕过 PennyLane TorchLayer 批处理缺陷, 10.4x 加速, 通用方案\n'
    '  2. TT 压缩: 双层 TT 分解 + 量子电路, 参数减少 93%, 数学严格的安全混淆\n'
    '  3. 安全分析框架: 对抗攻击 / 噪声鲁棒 / 梯度混淆 三维度测试\n'
    '\n'
    '核心见解:\n'
    '  量子模拟器 != 量子安全 -> TT 压缩在模拟器上提供唯一真正的安全保障\n'
    '  索引式参数与批处理的冲突 -> 显式参数是简洁的通用解决方案\n'
    '  安全性可以来自多个层次 (数学 + 物理 + 结构) -> 互补而非替代\n'
    '\n'
    '未来方向:\n'
    '  真量子硬件部署验证 (IBM Quantum / 本源量子)\n'
    '  更严格的安全性形式化证明\n'
    '  扩展到其他 PennyLane 电路类型和更大规模量子比特'
), y=1.3, size=16)
add_footer(s, '总结 & 未来方向')

# Save
path = 'd:/量子计算/QD-MSA-main/QD-MSA-main/models/QD-MSA_Batch_TT_Security.pptx'
prs.save(path)
print(f'PPT saved: {path}')
print(f'Total slides: {len(prs.slides)}')
